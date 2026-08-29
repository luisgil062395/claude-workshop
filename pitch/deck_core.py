#!/usr/bin/env python3
"""
Motor de render del pitch deck de Suma.

Una sola especificacion de formas se emite dos veces:
  1. PPTX editable (python-pptx) — el entregable.
  2. PNG de previsualizacion (Pillow) — para revisar el deck como conjunto.

Asi lo que se verifica es exactamente lo que se entrega.

Tokens tomados de src/styles/tokens.css (Suma Design System v1.0).
Nunca se define un color fuera de la tabla TOKENS.
"""

from __future__ import annotations
import os
import re
from dataclasses import dataclass, field

# ---------------------------------------------------------------- tokens

INK          = "131A2A"  # action-primary / text-primary
INK_HOVER    = "26313F"
SUBTLE       = "F1F2EF"  # action-subtle
POSITIVE     = "0E7A55"
AI           = "4A2BB5"  # ai-accent
AI_SUBTLE    = "EDE8FD"
DISABLED     = "B7BDC6"
BG           = "FBFBFA"
SURFACE      = "FFFFFF"
SURF_SUBTLE  = "F4F5F3"
BORDER       = "E4E6E1"
BORDER_SUB   = "EFF0EC"
T1           = "131A2A"
T2           = "5A6472"
T3           = "767F8C"
WHITE        = "FFFFFF"
SUCCESS      = "0E7A55"
WARNING      = "9A5B00"
ERROR        = "B3261E"
INFO         = "2A5CB8"
BRAND_GREEN  = "2E9E6B"
BRAND_DEEP   = "12314F"
BRAND_VIOLET = "5B34D6"
CHART = ["0B5C41", "9A5B00", "4A2BB5", "7FB8DE", "C4699B", "6B7280"]

# Escala tipografica del deck. Deriva de la del producto (una sola familia,
# pesos 400/600, tracking negativo en titulares) reescalada a 16:9.
DISPLAY, H1, H2, H3 = 44, 34, 24, 17
BODY_LG, BODY, BODY_SM = 15, 12.5, 11
LABEL, CAPTION = 9.5, 8.5
AMOUNT_LG, AMOUNT = 30, 16

SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.95

FONT_DIR = os.path.expanduser("~/Library/Fonts")
FONT_FILES = {"r": "Inter-Regular.ttf", "b": "Inter-Bold.ttf"}
FONT_NAME = "Inter"


# ---------------------------------------------------------------- formas

@dataclass
class Rect:
    x: float; y: float; w: float; h: float
    fill: str | None = None
    line: str | None = None
    lw: float = 1.0
    rad: float = 0.0
    dash: bool = False
    kind: str = "rect"          # rect | oval | tri


@dataclass
class Grad:
    x: float; y: float; w: float; h: float
    rad: float = 0.0
    ang: int = 0                # grados, 0 = izquierda -> derecha


@dataclass
class Img:
    path: str
    x: float; y: float; w: float; h: float


@dataclass
class Line:
    x1: float; y1: float; x2: float; y2: float
    color: str = BORDER
    lw: float = 1.0


@dataclass
class Text:
    x: float; y: float; w: float; h: float
    spans: object                # str | list[(text, weight, color)]
    size: float = BODY
    weight: str = "r"
    color: str = T1
    align: str = "l"             # l | c | r
    lh: float | None = None      # interlineado exacto en pt
    track: float = 0.0           # tracking en % del tamano
    lines: list = field(default_factory=list, repr=False)


def norm_spans(spans, weight, color):
    if isinstance(spans, str):
        return [(spans, weight, color)]
    return [(t, w or weight, c or color) for (t, w, c) in spans]


# ---------------------------------------------------------------- metricas

_font_cache: dict = {}


def pil_font(weight: str, size_pt: float, px_per_in: float):
    from PIL import ImageFont
    key = (weight, round(size_pt, 2), px_per_in)
    if key not in _font_cache:
        px = size_pt / 72.0 * px_per_in
        _font_cache[key] = ImageFont.truetype(
            os.path.join(FONT_DIR, FONT_FILES[weight]), int(round(px)))
    return _font_cache[key]


def text_width_pt(txt: str, weight: str, size: float, track: float) -> float:
    """Ancho en puntos, medido con la Inter real y con tracking aplicado."""
    f = pil_font(weight, size, 720.0)          # 720 px/in => 10 px por punto
    w = f.getlength(txt) / 10.0
    if track and len(txt) > 1:
        w += size * track / 100.0 * (len(txt) - 1)
    return w


def wrap(t: Text, safety: float = 0.94) -> list:
    """Parte los spans en lineas que caben en t.w. Devuelve [[(txt,w,c)]].

    Se trabaja sobre el texto concatenado con un atributo por caracter, de modo
    que un cambio de peso a mitad de palabra (p. ej. **antes**.) no inventa un
    espacio. El resultado se emite como saltos de linea explicitos en el PPTX,
    asi PowerPoint no puede re-partir el texto de otra forma.
    """
    limit = t.w * 72.0 * safety
    spans = norm_spans(t.spans, t.weight, t.color)
    full = "".join(x[0] for x in spans)
    attr = []
    for txt, w, c in spans:
        attr.extend([(w, c)] * len(txt))

    def seg(a, b):
        out, i = [], a
        while i < b:
            wc = attr[i]
            j = i
            while j < b and attr[j] == wc:
                j += 1
            out.append((full[i:j], wc[0], wc[1]))
            i = j
        return out

    lines, cur, cur_w = [], [], 0.0

    def push(segs):
        for txt, w, c in segs:
            if cur and cur[-1][1] == w and cur[-1][2] == c:
                cur[-1] = (cur[-1][0] + txt, w, c)
            else:
                cur.append((txt, w, c))

    pos = 0
    for para in full.split("\n"):
        for m in re.finditer(r"\S+", para):
            segs = seg(pos + m.start(), pos + m.end())
            ww = sum(text_width_pt(x, w, t.size, t.track) for x, w, _ in segs)
            gap = text_width_pt(" ", t.weight, t.size, t.track) if cur else 0.0
            if cur and cur_w + gap + ww > limit:
                lines.append(cur)
                cur, cur_w, gap = [], 0.0, 0.0
            if gap:
                push([(" ", cur[-1][1], cur[-1][2])])
            push(segs)
            cur_w += gap + ww
        lines.append(cur)
        cur, cur_w = [], 0.0
        pos += len(para) + 1

    lines = [l for l in lines if l]
    return lines or [[("", t.weight, t.color)]]


def nlines(text, w, size, weight="r", track=0.0) -> int:
    """Cuantas lineas ocupa un texto en un ancho dado. Para layouts que fluyen."""
    return len(wrap(Text(0, 0, w, 1.0, text, size=size, weight=weight,
                         track=track)))


def measure(shapes, slack: float = 0.12) -> list:
    """Calcula el wrap de cada Text y ajusta la caja a lo que realmente ocupa.

    Crecer la caja evita el recorte; lo que se reporta son los casos en que la
    caja crecio tanto que puede invadir al elemento siguiente y hay que mirar
    el render.
    """
    problems = []
    for s in shapes:
        if not isinstance(s, Text):
            continue
        s.lines = wrap(s)
        lh = s.lh or s.size * 1.35
        need = (len(s.lines) * lh) / 72.0
        if need > s.h + slack:
            problems.append((s, need))
        s.h = max(s.h, need)
    return problems


# ---------------------------------------------------------------- PPTX

def emit_pptx(slides, out_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.oxml.ns import qn
    import copy
    from lxml import etree

    def rgb(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W * 914400))
    prs.slide_height = Emu(int(SLIDE_H * 914400))
    blank = prs.slide_layouts[6]

    ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
    SHAPE = {"rect": MSO_SHAPE.ROUNDED_RECTANGLE, "oval": MSO_SHAPE.OVAL,
             "tri": MSO_SHAPE.ISOSCELES_TRIANGLE}

    for spec in slides:
        sl = prs.slides.add_slide(blank)

        for s in spec["shapes"]:
            if isinstance(s, (Rect, Grad)):
                kind = getattr(s, "kind", "rect")
                auto = SHAPE[kind] if not (kind == "rect" and not s.rad) \
                    else MSO_SHAPE.RECTANGLE
                sh = sl.shapes.add_shape(
                    auto, Inches(s.x), Inches(s.y), Inches(s.w), Inches(s.h))
                sh.shadow.inherit = False
                if kind == "rect" and s.rad:
                    sh.adjustments[0] = min(0.5, s.rad / min(s.w, s.h))

                if isinstance(s, Grad):
                    spPr = sh._element.spPr
                    for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
                        el = spPr.find(qn(tag))
                        if el is not None:
                            spPr.remove(el)
                    xml = (
                        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org'
                        '/drawingml/2006/main" rotWithShape="1"><a:gsLst>'
                        f'<a:gs pos="0"><a:srgbClr val="{BRAND_GREEN}"/></a:gs>'
                        f'<a:gs pos="52000"><a:srgbClr val="{BRAND_DEEP}"/></a:gs>'
                        f'<a:gs pos="100000"><a:srgbClr val="{BRAND_VIOLET}"/></a:gs>'
                        f'</a:gsLst><a:lin ang="{int(s.ang * 60000)}" scaled="0"/>'
                        '</a:gradFill>')
                    ln = spPr.find(qn("a:ln"))
                    node = etree.fromstring(xml)
                    spPr.insert(list(spPr).index(ln) if ln is not None
                                else len(spPr), node)
                    sh.line.fill.background()
                else:
                    if s.fill:
                        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(s.fill)
                    else:
                        sh.fill.background()
                    if s.line:
                        sh.line.color.rgb = rgb(s.line)
                        sh.line.width = Pt(s.lw)
                        if s.dash:
                            sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                    else:
                        sh.line.fill.background()
                sh.text_frame.text = ""

            elif isinstance(s, Line):
                cx = sl.shapes.add_connector(
                    1, Inches(s.x1), Inches(s.y1), Inches(s.x2), Inches(s.y2))
                cx.line.color.rgb = rgb(s.color)
                cx.line.width = Pt(s.lw)

            elif isinstance(s, Img):
                sl.shapes.add_picture(s.path, Inches(s.x), Inches(s.y),
                                      Inches(s.w), Inches(s.h))

            elif isinstance(s, Text):
                box = sl.shapes.add_textbox(
                    Inches(s.x), Inches(s.y), Inches(s.w), Inches(s.h))
                tf = box.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.TOP
                tf.margin_left = tf.margin_right = 0
                tf.margin_top = tf.margin_bottom = 0
                lh = s.lh or s.size * 1.35
                lines = s.lines or wrap(s)
                for i, line in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.alignment = ALIGN[s.align]
                    p.line_spacing = Pt(lh)
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    for txt, w, c in line:
                        r = p.add_run(); r.text = txt
                        r.font.name = FONT_NAME
                        r.font.size = Pt(s.size)
                        r.font.bold = (w == "b")
                        r.font.color.rgb = rgb(c)
                        if s.track:
                            r.font._rPr.set(
                                "spc", str(int(round(s.size * s.track))))

    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------- PNG

def emit_png(slides, out_dir, px=120):
    from PIL import Image, ImageDraw
    os.makedirs(out_dir, exist_ok=True)
    W, H = int(SLIDE_W * px), int(SLIDE_H * px)
    paths = []

    def C(h):
        return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))

    for n, spec in enumerate(slides, 1):
        im = Image.new("RGB", (W, H), C(BG))
        d = ImageDraw.Draw(im, "RGBA")

        for s in spec["shapes"]:
            if isinstance(s, (Rect, Grad)):
                x0, y0 = s.x * px, s.y * px
                x1, y1 = (s.x + s.w) * px, (s.y + s.h) * px
                if isinstance(s, Grad):
                    for i in range(int(x1 - x0)):
                        t = i / max(1, (x1 - x0) - 1)
                        if t < .52:
                            k = t / .52
                            a, b = C(BRAND_GREEN), C(BRAND_DEEP)
                        else:
                            k = (t - .52) / .48
                            a, b = C(BRAND_DEEP), C(BRAND_VIOLET)
                        col = tuple(int(a[j] + (b[j] - a[j]) * k) for j in range(3))
                        d.rectangle([x0 + i, y0, x0 + i + 1, y1], fill=col)
                    continue
                kind = getattr(s, "kind", "rect")
                fill = C(s.fill) if s.fill else None
                out = C(s.line) if s.line else None
                lw = max(1, int(round(s.lw * px / 72)))
                if kind == "oval":
                    d.ellipse([x0, y0, x1, y1], fill=fill, outline=out, width=lw)
                elif kind == "tri":
                    d.polygon([((x0 + x1) / 2, y0), (x1, y1), (x0, y1)],
                              fill=fill, outline=out)
                elif s.rad:
                    d.rounded_rectangle([x0, y0, x1, y1], radius=s.rad * px,
                                        fill=fill, outline=out, width=lw)
                else:
                    d.rectangle([x0, y0, x1, y1], fill=fill, outline=out, width=lw)

            elif isinstance(s, Line):
                d.line([s.x1 * px, s.y1 * px, s.x2 * px, s.y2 * px],
                       fill=C(s.color), width=max(1, int(round(s.lw * px / 72))))

            elif isinstance(s, Img):
                sub = Image.open(s.path).convert("RGBA").resize(
                    (max(1, int(s.w * px)), max(1, int(s.h * px))), Image.LANCZOS)
                im.paste(sub, (int(s.x * px), int(s.y * px)), sub)

            elif isinstance(s, Text):
                lines = s.lines or wrap(s)
                lh = (s.lh or s.size * 1.35) / 72 * px
                asc = s.size / 72 * px
                for i, line in enumerate(lines):
                    total = sum(text_width_pt(t, w, s.size, s.track)
                                for t, w, _ in line) / 72 * px
                    if s.align == "c":
                        cx = s.x * px + (s.w * px - total) / 2
                    elif s.align == "r":
                        cx = s.x * px + s.w * px - total
                    else:
                        cx = s.x * px
                    cy = s.y * px + i * lh + (lh - asc) * 0.42
                    for t, w, c in line:
                        f = pil_font(w, s.size, px)
                        if s.track:
                            step = s.size * s.track / 100 / 72 * px
                            for ch in t:
                                d.text((cx, cy), ch, font=f, fill=C(c))
                                cx += f.getlength(ch) + step
                        else:
                            d.text((cx, cy), t, font=f, fill=C(c))
                            cx += f.getlength(t)

        p = os.path.join(out_dir, f"slide-{n:02d}.png")
        im.save(p); paths.append(p)
    return paths
